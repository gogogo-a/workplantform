"""
各种文件格式的提取器实现
"""
from typing import List
from pathlib import Path
from .base_extractor import BaseExtractor
from log import logger


class TextExtractor(BaseExtractor):
    """纯文本文件提取器（.txt, .md, .log 等）"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.txt', '.md', '.log', '.text']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            logger.info(f"✓ 文本文件提取成功: {Path(file_path).name}, 长度: {len(content)}")
            return content
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    content = f.read()
                logger.info(f"✓ 文本文件提取成功（GBK编码）: {Path(file_path).name}")
                return content
            except Exception as e:
                logger.error(f"✗ 文本文件提取失败: {e}")
                raise


class PDFExtractor(BaseExtractor):
    """PDF 文件提取器（使用 PyMuPDF，支持图片和表格，带 OCR 识别）"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.pdf']
    
    def _ocr_image(self, image_bytes: bytes) -> str:
        """
        使用 OCR 识别图片中的文字
        
        Args:
            image_bytes: 图片字节流
            
        Returns:
            str: 识别出的文字，失败返回空字符串
        """
        try:
            from PIL import Image
            import pytesseract
            import io
            
            # 将字节流转为图片对象
            image = Image.open(io.BytesIO(image_bytes))
            
            # OCR 识别（支持中英文）
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            
            return text.strip()
            
        except ImportError:
            logger.debug("pytesseract 或 Pillow 未安装，跳过 OCR 识别")
            return ""
        except Exception as e:
            logger.debug(f"OCR 识别失败: {e}")
            return ""
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            content_parts = []
            image_count = 0
            ocr_success_count = 0
            table_count = 0
            total_pages = len(doc)  # 🔥 先保存页数，避免关闭后访问
            
            for page_num in range(total_pages):
                page = doc[page_num]
                
                # 提取文本
                text = page.get_text("text")
                
                # 提取图片并 OCR 识别
                images = page.get_images()
                if images:
                    image_count += len(images)
                    text += f"\n\n--- 页面 {page_num + 1} 的图片内容 ({len(images)} 张) ---"
                    
                    for img_index, img in enumerate(images):
                        try:
                            # 获取图片引用
                            xref = img[0]
                            
                            # 提取图片字节流
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            # OCR 识别
                            ocr_text = self._ocr_image(image_bytes)
                            
                            if ocr_text:
                                text += f"\n\n[图片 {img_index + 1} 识别内容]:\n{ocr_text}"
                                ocr_success_count += 1
                            else:
                                text += f"\n[图片 {img_index + 1}: 无法识别文字]"
                                
                        except Exception as e:
                            logger.debug(f"处理图片失败 (页 {page_num + 1}, 图 {img_index + 1}): {e}")
                            text += f"\n[图片 {img_index + 1}: 提取失败]"
                
                # 提取表格（使用文本块检测）
                blocks = page.get_text("dict")["blocks"]
                table_blocks = [b for b in blocks if b.get("type") == 0 and len(b.get("lines", [])) > 3]
                if table_blocks:
                    table_count += len(table_blocks)
                    text += f"\n\n[页面 {page_num + 1} 包含 {len(table_blocks)} 个表格结构]"
                
                if text.strip():
                    content_parts.append(f"--- 第 {page_num + 1} 页 ---\n{text}")
            
            doc.close()
            
            content = "\n\n".join(content_parts)
            logger.info(
                f"✓ PDF提取成功: {Path(file_path).name}, "
                f"页数: {total_pages}, 长度: {len(content)}, "
                f"图片: {image_count} (OCR成功: {ocr_success_count}), "
                f"表格结构: {table_count}"
            )
            return content
            
        except ImportError as e:
            # 降级到 PyPDFLoader
            logger.warning(f"PyMuPDF 未安装，降级使用 PyPDFLoader（不支持图片提取）: {e}")
            try:
                from langchain_community.document_loaders import PyPDFLoader
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                content = "\n\n".join([doc.page_content for doc in documents])
                logger.info(f"✓ PDF提取成功（PyPDFLoader）: {Path(file_path).name}, 页数: {len(documents)}")
                return content
            except Exception as e:
                logger.error(f"✗ PDF提取失败: {e}")
                raise
        except Exception as e:
            logger.error(f"✗ PDF提取失败: {e}")
            raise


class WordExtractor(BaseExtractor):
    """Word 文档提取器（.docx, .doc，支持图片 OCR 识别）"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.docx', '.doc']
    
    def _convert_doc_to_docx(self, doc_path: str) -> str:
        """
        使用 LibreOffice 将 .doc 转换为 .docx
        
        Args:
            doc_path: .doc 文件路径
            
        Returns:
            str: 转换后的 .docx 文件路径
            
        Raises:
            Exception: 转换失败时抛出异常
        """
        import subprocess
        import tempfile
        import shutil
        
        doc_path_obj = Path(doc_path)
        
        # 创建临时目录存放转换后的文件
        temp_dir = tempfile.mkdtemp(prefix="doc_convert_")
        
        try:
            # 执行 LibreOffice 转换命令
            cmd = [
                'soffice',  # LibreOffice 命令（macOS/Linux）
                '--headless',  # 无界面模式
                '--convert-to', 'docx',  # 转换为 docx
                '--outdir', temp_dir,  # 输出目录
                doc_path
            ]
            
            logger.info(f"开始转换 .doc → .docx: {doc_path_obj.name}")
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=30  # 30秒超时
            )
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice 转换失败: {result.stderr}")
            
            # 转换后的文件路径
            docx_path = Path(temp_dir) / f"{doc_path_obj.stem}.docx"
            
            if not docx_path.exists():
                raise Exception(f"转换后的文件不存在: {docx_path}")
            
            logger.info(f"✓ .doc 转换成功: {doc_path_obj.name} → {docx_path.name}")
            return str(docx_path)
            
        except FileNotFoundError:
            raise Exception(
                "LibreOffice 未安装。请安装 LibreOffice:\n"
                "  macOS: brew install --cask libreoffice\n"
                "  Ubuntu: sudo apt-get install libreoffice\n"
                "  CentOS: sudo yum install libreoffice"
            )
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice 转换超时（30秒）")
        except Exception as e:
            # 清理临时目录
            if Path(temp_dir).exists():
                shutil.rmtree(temp_dir, ignore_errors=True)
            raise
    
    def _ocr_image(self, image_bytes: bytes) -> str:
        """
        使用 OCR 识别图片中的文字
        
        Args:
            image_bytes: 图片字节流
            
        Returns:
            str: 识别出的文字，失败返回空字符串
        """
        try:
            from PIL import Image
            import pytesseract
            import io
            
            # 将字节流转为图片对象
            image = Image.open(io.BytesIO(image_bytes))
            
            # OCR 识别（支持中英文）
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            
            return text.strip()
            
        except ImportError:
            logger.debug("pytesseract 或 Pillow 未安装，跳过 OCR 识别")
            return ""
        except Exception as e:
            logger.debug(f"OCR 识别失败: {e}")
            return ""
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        file_ext = Path(file_path).suffix.lower()
        converted_file = None  # 用于清理转换后的临时文件
        
        # 🔥 .doc 格式：先用 LibreOffice 转换为 .docx，然后支持图片 OCR
        if file_ext == '.doc':
            logger.info(f"检测到 .doc 格式（旧版 Word），尝试使用 LibreOffice 转换...")
            
            try:
                # 方案 1：LibreOffice 转换（支持图片 OCR）
                import shutil
                
                converted_file = self._convert_doc_to_docx(file_path)
                logger.info(f"使用转换后的 .docx 文件进行提取（支持图片 OCR）")
                
                # 递归调用，处理转换后的 .docx
                content = self.extract_from_file(converted_file)
                
                # 清理转换后的临时文件和目录
                temp_dir = Path(converted_file).parent
                if temp_dir.exists() and "doc_convert_" in temp_dir.name:
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    logger.debug(f"已清理临时目录: {temp_dir}")
                
                return content
                
            except Exception as convert_error:
                # 方案 2：LibreOffice 转换失败，降级使用 Docx2txtLoader（不支持图片 OCR）
                logger.warning(f"LibreOffice 转换失败: {convert_error}")
                logger.info(f"降级使用 Docx2txtLoader 提取（不支持图片 OCR）")
                
                try:
                    from langchain_community.document_loaders import Docx2txtLoader
                    loader = Docx2txtLoader(file_path)
                    documents = loader.load()
                    content = "\n\n".join([doc.page_content for doc in documents])
                    logger.info(f"✓ Word文档提取成功（.doc 降级模式）: {Path(file_path).name}, 长度: {len(content)}")
                    return content
                except Exception as e:
                    logger.error(f"✗ .doc 文档提取失败（降级模式也失败）: {e}")
                    raise
            finally:
                # 确保清理临时文件
                if converted_file and Path(converted_file).exists():
                    import shutil
                    temp_dir = Path(converted_file).parent
                    if temp_dir.exists() and "doc_convert_" in temp_dir.name:
                        shutil.rmtree(temp_dir, ignore_errors=True)
        
        # 🔥 .docx 格式使用 python-docx（支持图片 OCR）
        try:
            from docx import Document
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
            
            doc = Document(file_path)
            content_parts = []
            image_count = 0
            ocr_success_count = 0
            
            # 遍历文档的所有元素（段落、表格、图片）
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    # 段落
                    paragraph = Paragraph(element, doc)
                    text = paragraph.text.strip()
                    
                    # 检查段落中是否有图片
                    for run in paragraph.runs:
                        # 获取图片
                        for inline_shape in run.element.xpath('.//a:blip'):
                            try:
                                image_count += 1
                                # 获取图片关系 ID
                                rId = inline_shape.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                
                                if rId:
                                    # 获取图片数据
                                    image_part = doc.part.related_parts[rId]
                                    image_bytes = image_part.blob
                                    
                                    # OCR 识别
                                    ocr_text = self._ocr_image(image_bytes)
                                    
                                    if ocr_text:
                                        text += f"\n\n[图片 {image_count} 识别内容]:\n{ocr_text}"
                                        ocr_success_count += 1
                                    else:
                                        text += f"\n[图片 {image_count}: 无法识别文字]"
                            except Exception as e:
                                logger.debug(f"处理 Word 图片失败: {e}")
                                text += f"\n[图片 {image_count}: 提取失败]"
                    
                    if text:
                        content_parts.append(text)
                
                elif isinstance(element, CT_Tbl):
                    # 表格
                    table = Table(element, doc)
                    table_text = "\n=== 表格 ===\n"
                    
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        table_text += row_text + "\n"
                    
                    content_parts.append(table_text)
            
            content = "\n\n".join(content_parts)
            logger.info(
                f"✓ Word文档提取成功（.docx 格式）: {Path(file_path).name}, "
                f"长度: {len(content)}, "
                f"图片: {image_count} (OCR成功: {ocr_success_count})"
            )
            return content
            
        except ImportError:
            # python-docx 未安装，降级到 Docx2txtLoader
            logger.warning("python-docx 未安装，降级使用 Docx2txtLoader（不支持图片 OCR）")
            try:
                from langchain_community.document_loaders import Docx2txtLoader
                loader = Docx2txtLoader(file_path)
                documents = loader.load()
                content = "\n\n".join([doc.page_content for doc in documents])
                logger.info(f"✓ Word文档提取成功（Docx2txtLoader）: {Path(file_path).name}, 长度: {len(content)}")
                return content
            except Exception as e:
                logger.error(f"✗ Word文档提取失败: {e}")
                raise
        except Exception as e:
            # .docx 文件处理失败，可能是文件损坏
            logger.error(f"✗ .docx 文档提取失败: {e}")
            logger.warning("尝试使用 Docx2txtLoader 降级处理...")
            try:
                from langchain_community.document_loaders import Docx2txtLoader
                loader = Docx2txtLoader(file_path)
                documents = loader.load()
                content = "\n\n".join([doc.page_content for doc in documents])
                logger.info(f"✓ Word文档提取成功（降级模式）: {Path(file_path).name}, 长度: {len(content)}")
                return content
            except Exception as fallback_error:
                logger.error(f"✗ Word文档提取失败（降级模式也失败）: {fallback_error}")
                raise


class PowerPointExtractor(BaseExtractor):
    """PowerPoint 提取器（.pptx, .ppt，支持图片 OCR 识别）"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.pptx', '.ppt']
    
    def _convert_ppt_to_pptx(self, ppt_path: str) -> str:
        """
        使用 LibreOffice 将 .ppt 转换为 .pptx
        
        Args:
            ppt_path: .ppt 文件路径
            
        Returns:
            str: 转换后的 .pptx 文件路径
            
        Raises:
            Exception: 转换失败时抛出异常
        """
        import subprocess
        import tempfile
        import shutil
        
        ppt_path_obj = Path(ppt_path)
        
        # 创建临时目录存放转换后的文件
        temp_dir = tempfile.mkdtemp(prefix="ppt_convert_")
        
        try:
            # 执行 LibreOffice 转换命令
            cmd = [
                'soffice',  # LibreOffice 命令（macOS/Linux）
                '--headless',  # 无界面模式
                '--convert-to', 'pptx',  # 转换为 pptx
                '--outdir', temp_dir,  # 输出目录
                ppt_path
            ]
            
            logger.info(f"开始转换 .ppt → .pptx: {ppt_path_obj.name}")
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60  # 60秒超时（PPT 文件可能较大）
            )
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice 转换失败: {result.stderr}")
            
            # 转换后的文件路径
            pptx_path = Path(temp_dir) / f"{ppt_path_obj.stem}.pptx"
            
            if not pptx_path.exists():
                raise Exception(f"转换后的文件不存在: {pptx_path}")
            
            logger.info(f"✓ .ppt 转换成功: {ppt_path_obj.name} → {pptx_path.name}")
            return str(pptx_path)
            
        except FileNotFoundError:
            raise Exception(
                "LibreOffice 未安装。请安装 LibreOffice:\n"
                "  macOS: brew install --cask libreoffice\n"
                "  Ubuntu: sudo apt-get install libreoffice\n"
                "  CentOS: sudo yum install libreoffice"
            )
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice 转换超时（60秒）")
        except Exception as e:
            # 清理临时目录
            if Path(temp_dir).exists():
                shutil.rmtree(temp_dir, ignore_errors=True)
            raise
    
    def _ocr_image(self, image_bytes: bytes) -> str:
        """
        使用 OCR 识别图片中的文字
        
        Args:
            image_bytes: 图片字节流
            
        Returns:
            str: 识别出的文字，失败返回空字符串
        """
        try:
            from PIL import Image
            import pytesseract
            import io
            
            # 将字节流转为图片对象
            image = Image.open(io.BytesIO(image_bytes))
            
            # OCR 识别（支持中英文）
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            
            return text.strip()
            
        except ImportError:
            logger.debug("pytesseract 或 Pillow 未安装，跳过 OCR 识别")
            return ""
        except Exception as e:
            logger.debug(f"OCR 识别失败: {e}")
            return ""
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        file_ext = Path(file_path).suffix.lower()
        converted_file = None  # 用于清理转换后的临时文件
        
        # 🔥 .ppt 格式：先用 LibreOffice 转换为 .pptx，然后支持图片 OCR
        if file_ext == '.ppt':
            logger.info(f"检测到 .ppt 格式（旧版 PowerPoint），尝试使用 LibreOffice 转换...")
            
            try:
                # LibreOffice 转换（支持图片 OCR）
                import shutil
                
                converted_file = self._convert_ppt_to_pptx(file_path)
                logger.info(f"使用转换后的 .pptx 文件进行提取（支持图片 OCR）")
                
                # 递归调用，处理转换后的 .pptx
                content = self.extract_from_file(converted_file)
                
                # 清理转换后的临时文件和目录
                temp_dir = Path(converted_file).parent
                if temp_dir.exists() and "ppt_convert_" in temp_dir.name:
                    shutil.rmtree(temp_dir, ignore_errors=True)
                    logger.debug(f"已清理临时目录: {temp_dir}")
                
                return content
                
            except Exception as convert_error:
                # LibreOffice 转换失败，给出提示
                logger.error(f"LibreOffice 转换失败: {convert_error}")
                logger.warning(
                    f"⚠️ .ppt 格式自动转换失败。\n"
                    f"建议：请手动将文件转换为 .pptx 格式后重新上传。\n"
                    f"转换方法：在 PowerPoint 中打开文件，另存为 .pptx 格式。"
                )
                raise ValueError(
                    f"不支持 .ppt 格式（自动转换失败）。"
                    f"请将文件手动转换为 .pptx 格式后重新上传。原因: {str(convert_error)}"
                )
            finally:
                # 确保清理临时文件
                if converted_file and Path(converted_file).exists():
                    import shutil
                    temp_dir = Path(converted_file).parent
                    if temp_dir.exists() and "ppt_convert_" in temp_dir.name:
                        shutil.rmtree(temp_dir, ignore_errors=True)
        
        # 🔥 .pptx 格式使用 python-pptx（支持图片 OCR）
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            prs = Presentation(file_path)
            content_parts = []
            total_images = 0
            total_ocr_success = 0
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {i} ---")
                slide_images = 0
                
                for shape in slide.shapes:
                    # 提取文本
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # 提取图片
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            slide_images += 1
                            total_images += 1
                            
                            # 获取图片数据
                            image = shape.image
                            image_bytes = image.blob
                            
                            # OCR 识别
                            ocr_text = self._ocr_image(image_bytes)
                            
                            if ocr_text:
                                slide_text.append(f"\n[图片 {slide_images} 识别内容]:\n{ocr_text}")
                                total_ocr_success += 1
                            else:
                                slide_text.append(f"[图片 {slide_images}: 无法识别文字]")
                        except Exception as e:
                            logger.debug(f"处理 PPT 图片失败 (Slide {i}): {e}")
                            slide_text.append(f"[图片 {slide_images}: 提取失败]")
                
                if len(slide_text) > 1:  # 有内容
                    content_parts.append("\n".join(slide_text))
            
            content = "\n\n".join(content_parts)
            logger.info(
                f"✓ PPT提取成功（.pptx 格式）: {Path(file_path).name}, "
                f"幻灯片数: {len(prs.slides)}, 长度: {len(content)}, "
                f"图片: {total_images} (OCR成功: {total_ocr_success})"
            )
            return content
        except Exception as e:
            logger.error(f"✗ .pptx 文档提取失败: {e}")
            raise


class ExcelExtractor(BaseExtractor):
    """Excel 提取器（.xlsx, .xls）"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.xlsx', '.xls']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            import pandas as pd
            
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # 转换为文本
                sheet_text = f"=== Sheet: {sheet_name} ===\n"
                sheet_text += df.to_string(index=False)
                content_parts.append(sheet_text)
            
            content = "\n\n".join(content_parts)
            logger.info(f"✓ Excel提取成功: {Path(file_path).name}, 工作表数: {len(excel_file.sheet_names)}, 长度: {len(content)}")
            return content
        except Exception as e:
            logger.error(f"✗ Excel提取失败: {e}")
            raise


class CSVExtractor(BaseExtractor):
    """CSV 提取器"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.csv']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            import pandas as pd
            
            # 尝试不同编码
            for encoding in ['utf-8', 'gbk', 'gb2312', 'latin1']:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    content = df.to_string(index=False)
                    logger.info(f"✓ CSV提取成功: {Path(file_path).name}, 行数: {len(df)}, 编码: {encoding}")
                    return content
                except UnicodeDecodeError:
                    continue
            
            raise ValueError("无法识别CSV文件编码")
        except Exception as e:
            logger.error(f"✗ CSV提取失败: {e}")
            raise


class HTMLExtractor(BaseExtractor):
    """HTML 提取器"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.html', '.htm']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 移除 script 和 style 标签
            for script in soup(["script", "style"]):
                script.decompose()
            
            # 获取文本
            text = soup.get_text()
            
            # 清理多余空白
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            logger.info(f"✓ HTML提取成功: {Path(file_path).name}, 长度: {len(text)}")
            return text
        except Exception as e:
            logger.error(f"✗ HTML提取失败: {e}")
            raise


class RTFExtractor(BaseExtractor):
    """RTF 富文本提取器"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.rtf']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            from striprtf.striprtf import rtf_to_text
            
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            logger.info(f"✓ RTF提取成功: {Path(file_path).name}, 长度: {len(text)}")
            return text
        except Exception as e:
            logger.error(f"✗ RTF提取失败: {e}")
            raise


class EPUBExtractor(BaseExtractor):
    """EPUB 电子书提取器"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.epub']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(file_path)
            content_parts = []
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text()
                    if text.strip():
                        content_parts.append(text.strip())
            
            content = "\n\n".join(content_parts)
            logger.info(f"✓ EPUB提取成功: {Path(file_path).name}, 章节数: {len(content_parts)}, 长度: {len(content)}")
            return content
        except Exception as e:
            logger.error(f"✗ EPUB提取失败: {e}")
            raise


class JSONExtractor(BaseExtractor):
    """JSON 文件提取器"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.json']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # 转换为可读文本
            content = json.dumps(data, ensure_ascii=False, indent=2)
            logger.info(f"✓ JSON提取成功: {Path(file_path).name}, 长度: {len(content)}")
            return content
        except Exception as e:
            logger.error(f"✗ JSON提取失败: {e}")
            raise


class XMLExtractor(BaseExtractor):
    """XML 文件提取器"""
    
    def get_supported_extensions(self) -> List[str]:
        return ['.xml']
    
    def extract_from_file(self, file_path: str) -> str:
        self.validate_file(file_path)
        
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                xml_content = f.read()
            
            soup = BeautifulSoup(xml_content, 'xml')
            text = soup.get_text()
            
            # 清理多余空白
            lines = (line.strip() for line in text.splitlines())
            text = '\n'.join(line for line in lines if line)
            
            logger.info(f"✓ XML提取成功: {Path(file_path).name}, 长度: {len(text)}")
            return text
        except Exception as e:
            logger.error(f"✗ XML提取失败: {e}")
            raise

